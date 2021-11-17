# -*- encoding: utf-8 -*-
'''
@File   :   text_parsing.py
@Time   :   2021/04/20 17:00:00
@Author :   qjzhzw
@Email  :   450388261@qq.com
@Desc   :   解析各种格式的文件 (txt/docx/pptx/pdf/html/eml)
'''

import re
import docx
import email
from pptx import Presentation
from pptx.util import Cm, Pt
from pdfminer.converter import PDFPageAggregator
from pdfminer.layout import LAParams
from pdfminer.pdfparser import PDFParser, PDFDocument
from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
from bs4 import BeautifulSoup


def read_txt(in_path):
    """
    读取txt文件
    """
    with open(in_path, "r", encoding="utf-8") as f:
        data = [line.strip() for line in f if len(line.strip()) != 0]
    return data


def read_docx(in_path):
    """
    读取docx文件
    """
    data = []
    f = docx.Document(in_path)
    for para in f.paragraphs:
        text = para.text.strip()
        if len(text) != 0:
            data.append(text)
    return data


def read_pptx(in_path):
    """
    读取pptx文件
    """
    data = []
    prs = Presentation(in_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    data.append(paragraph.text)
    data = [d.strip() for d in data if len(d.strip()) != 0]
    return data


def read_pdf(in_path):
    """
    读取pdf文件
    """
    data = []

    f = open(in_path, "rb")
    parser = PDFParser(f)
    doc = PDFDocument()
    parser.set_document(doc)
    doc.set_parser(parser)
    doc.initialize("")
    resource = PDFResourceManager()
    laparam = LAParams()
    device = PDFPageAggregator(resource, laparams=laparam)
    interpreter = PDFPageInterpreter(resource, device)
    for page in doc.get_pages():
        interpreter.process_page(page)
        layout = device.get_result()
        for out in layout:
            data = out.get_text().split("\n")
    data = [d.strip() for d in data if len(d.strip()) != 0]
    f.close()
    return data


def read_html(in_path):
    """
    读取html文件
    """
    data = []
    f = open(in_path, "r", encoding="gbk")
    soup = BeautifulSoup(f.read(), "html.parser")
    data = soup.get_text()
    data = data.split("\n")
    data = [d.strip() for d in data if len(d.strip()) != 0 and \
            not re.search(r"p.p1", d) and \
            not re.search(r"span.s1", d) and \
            not re.search(r"span.s2", d)]
    f.close()
    return data


def read_eml(in_path):
    """
    读取eml文件
    """
    data = []
    f = open(in_path, "r", encoding="utf-8")
    msg = email.message_from_file(f)
    data = str(msg).split("\n")
    data = [d.strip() for d in data if len(d.strip()) != 0]
    f.close()
    return data

